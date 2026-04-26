# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "cairosvg",
#     "python-pptx",
# ]
# ///
"""Parametric isometric studded-brick generator.

Three perspectives per block:
  - iso:  30 degree isometric (matches original style), right- or left-facing
  - top:  straight-down view, square top + circular studs
  - side: front elevation, rectangle with stud bumps

Block defined as (W, D, H_units) where:
  W = studs wide along iso-right axis
  D = studs deep along iso-back-left axis
  H_units = height in plate-units (1 = plate, 3 = brick, 6 = tall brick)
"""

from __future__ import annotations
import argparse
import colorsys
import math
import os
import re
import shutil
from concurrent.futures import ProcessPoolExecutor
from dataclasses import dataclass
from pathlib import Path


VB_RE = re.compile(r'viewBox="([^"]+)"')


# ---------- Geometry ----------
# Real-world studded brick: stud pitch 8mm, brick body 9.6mm (3 plates), plate 3.2mm.
# PLATE_H=12 with U=30 gives brick/stud = 1.2 and plate/stud = 0.4 — matches real-world ratios.
U = 30                # screen pixels per stud along iso axis
S = U // 2            # iso vertical step per stud (2:1 pixel iso)
PLATE_H = 12          # screen height per plate-unit (brick = 3 units = 36px)
STUD_RX = 9           # stud ellipse radius x (~0.6 of cell width, matches real-world brick ratio)
STUD_RY = STUD_RX / 2 # iso-foreshortened y radius
STUD_RAISE = 4        # vertical screen offset between stud base and top
STROKE_W = 0.6


# ---------- Palette ----------
# For each colour we derive: top, left wall, right wall, stud top, stud side,
# outline. The mantel set's values are hand-tuned; classic is computed.
@dataclass(frozen=True)
class Palette:
    name: str
    top: str       # main top face
    left: str      # left visible wall (medium shade)
    right: str     # right visible wall (darker shade)
    stud_top: str  # peg top ellipse (lighter than top)
    stud_side: str # peg base ellipse (darker than top)
    outline: str   # stroke colour
    accent: bool = False  # only render on ACCENT_SLUGS (for light-on-light tones)


def _hex_to_rgb(hx: str) -> tuple[float, float, float]:
    hx = hx.lstrip("#")
    return (
        int(hx[0:2], 16) / 255,
        int(hx[2:4], 16) / 255,
        int(hx[4:6], 16) / 255,
    )


def _hls_to_hex(h: float, l: float, s: float) -> str:
    r, g, b = colorsys.hls_to_rgb(h, max(0.0, min(1.0, l)), max(0.0, min(1.0, s)))
    return f"#{int(round(r * 255)):02X}{int(round(g * 255)):02X}{int(round(b * 255)):02X}"


def auto_ramp(name: str, base: str, accent: bool = False) -> Palette:
    """Derive a 6-colour shade ramp from a single base hex.

    Hand-tuned ramps (like the mantel ones) look a little richer, but this is
    good enough for generic palettes and keeps each entry to one line.
    """
    r, g, b = _hex_to_rgb(base)
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    # Light bases (l > 0.75) can't be darkened into a readable outline, so
    # drop to a dark desaturated grey instead of a darker tint.
    if l > 0.75:
        outline = _hls_to_hex(h, 0.20, s * 0.5)
    else:
        outline = _hls_to_hex(h, max(0.05, l * 0.30), s)
    return Palette(
        name=name,
        top=base,
        left=_hls_to_hex(h, l - 0.08, s),
        right=_hls_to_hex(h, l - 0.16, s),
        stud_top=_hls_to_hex(h, l + 0.12, s),
        stud_side=_hls_to_hex(h, l - 0.08, s),
        outline=outline,
        accent=accent,
    )


# ---------- Palette sets ----------
# Each set is a dict of palette-name -> Palette. `--palettes` picks one or
# more sets; filenames embed the palette name so sets can coexist on disk.

MANTEL_PALETTES: dict[str, Palette] = {
    "ocean": Palette(
        name="ocean",
        top="#1E5E82", left="#174E6E", right="#103D57",
        stud_top="#3D89AF", stud_side="#174E6E", outline="#042031",
    ),
    "flamingo": Palette(
        name="flamingo",
        top="#D86E89", left="#C25079", right="#9E3456",
        stud_top="#E690A2", stud_side="#C25079", outline="#550E28",
    ),
    "sky": Palette(
        name="sky",
        top="#81CCEA", left="#59BAE2", right="#329DD3",
        stud_top="#B8E5F5", stud_side="#59BAE2", outline="#064260",
    ),
    "deep-ocean": Palette(
        name="deep-ocean",
        top="#002A41", left="#001E2F", right="#001421",
        stud_top="#073B58", stud_side="#001E2F", outline="#000810",
    ),
    "cloud": Palette(
        name="cloud",
        top="#EEF9FD", left="#D9F1F9", right="#B8E5F5",
        stud_top="#FFFFFF", stud_side="#D9F1F9", outline="#1E5E82",
        accent=True,
    ),
    # Cool-cast black/grey to sit alongside the ocean-biased mantel tones.
    "black": Palette(
        name="black",
        top="#141820", left="#0B0F16", right="#04070F",
        stud_top="#2B3240", stud_side="#0B0F16", outline="#000000",
    ),
    "grey": Palette(
        name="grey",
        top="#64707C", left="#515B66", right="#3E4650",
        stud_top="#8593A0", stud_side="#515B66", outline="#1A2028",
    ),
}

# Standard studded-brick reference hexes for the classic palette.
CLASSIC_PALETTES: dict[str, Palette] = {
    "l-red":    auto_ramp("l-red",    "#C91A09"),
    "l-blue":   auto_ramp("l-blue",   "#0055BF"),
    "l-yellow": auto_ramp("l-yellow", "#F2CD37"),
    "l-green":  auto_ramp("l-green",  "#237841"),
    "l-white":  auto_ramp("l-white",  "#F4F4F4"),
    "l-black":  auto_ramp("l-black",  "#1B2A34"),
    # "Light Bluish Gray" — the canonical modern studded-brick grey. No
    # separate black entry: l-black above already covers that end.
    "l-grey":   auto_ramp("l-grey",   "#A0A5A9"),
}

# Anthropic-aligned palette. Slate/Smoke/Ivory cover the greyscale range
# (Smoke is the medium grey from the brand's "Cloud" family, renamed to
# avoid clashing with mantel's cloud accent); Book Cloth / Kraft / Manilla
# are the signature warm trio. Ivory is marked accent — too light to read
# across the full catalogue on a transparent background.
ANTHROPIC_PALETTES: dict[str, Palette] = {
    "slate":      auto_ramp("slate",      "#262625"),
    "smoke":      auto_ramp("smoke",      "#91918D"),
    "ivory":      auto_ramp("ivory",      "#F0F0EB", accent=True),
    "book-cloth": auto_ramp("book-cloth", "#CC785C"),
    "kraft":      auto_ramp("kraft",      "#D4A27F"),
    "manilla":    auto_ramp("manilla",    "#EBDBBC"),
    # Warm-neutral black/grey that reads as a deliberate part of the
    # anthropic family rather than a cool intrusion. Slate (#262625) is
    # already near-black but warm-toned; pure black reads distinctly darker.
    # Grey slots between slate and smoke (#91918D) rather than duplicating
    # either.
    "black":      auto_ramp("black",      "#0E0E0D"),
    "grey":       auto_ramp("grey",       "#5A5A57"),
}

PALETTE_SETS: dict[str, dict[str, Palette]] = {
    "mantel": MANTEL_PALETTES,
    "classic": CLASSIC_PALETTES,
    "anthropic": ANTHROPIC_PALETTES,
}

# Accent palettes (Palette.accent=True) render only on ACCENT_SLUGS; full
# catalogue output would be near-invisible on a transparent background.
ACCENT_SLUGS = frozenset({"brick-2x2", "brick-2x4"})


# ---------- SVG helpers ----------
def fmt(x: float) -> str:
    """Trim trailing zeros from numeric output."""
    if abs(x - round(x)) < 1e-6:
        return str(int(round(x)))
    return f"{x:.3f}".rstrip("0").rstrip(".")


def poly(points: list[tuple[float, float]], fill: str, stroke: str) -> str:
    pts = " ".join(f"{fmt(x)},{fmt(y)}" for x, y in points)
    return (
        f'<polygon points="{pts}" fill="{fill}" '
        f'stroke="{stroke}" stroke-width="{STROKE_W}" stroke-linejoin="round"/>'
    )


def ellipse(cx: float, cy: float, rx: float, ry: float, fill: str, stroke: str) -> str:
    return (
        f'<ellipse cx="{fmt(cx)}" cy="{fmt(cy)}" rx="{fmt(rx)}" ry="{fmt(ry)}" '
        f'fill="{fill}" stroke="{stroke}" stroke-width="{STROKE_W}"/>'
    )


def circle(cx: float, cy: float, r: float, fill: str, stroke: str) -> str:
    return (
        f'<circle cx="{fmt(cx)}" cy="{fmt(cy)}" r="{fmt(r)}" '
        f'fill="{fill}" stroke="{stroke}" stroke-width="{STROKE_W}"/>'
    )


def rect(x: float, y: float, w: float, h: float, fill: str, stroke: str, rx: float = 0) -> str:
    extra = f' rx="{fmt(rx)}"' if rx else ""
    return (
        f'<rect x="{fmt(x)}" y="{fmt(y)}" width="{fmt(w)}" height="{fmt(h)}"'
        f'{extra} fill="{fill}" stroke="{stroke}" stroke-width="{STROKE_W}"/>'
    )


def _zigzag(
    start: tuple[float, float], end: tuple[float, float],
    n_peaks: int, amplitude: float,
) -> list[tuple[float, float]]:
    """Polyline from start to end with alternating perpendicular peaks.

    The first interior peak offsets in the CCW-perpendicular direction; peaks
    alternate from there. Used to draw the jagged snap edge for broken bricks.
    """
    ax, ay = start
    bx, by = end
    dx, dy = bx - ax, by - ay
    length = (dx * dx + dy * dy) ** 0.5 or 1.0
    px, py = -dy / length, dx / length
    points: list[tuple[float, float]] = [start]
    for i in range(1, n_peaks + 1):
        t = i / (n_peaks + 1)
        sign = 1 if i % 2 == 1 else -1
        points.append((
            ax + dx * t + sign * amplitude * px,
            ay + dy * t + sign * amplitude * py,
        ))
    points.append(end)
    return points


def _polyline_path(points: list[tuple[float, float]], closed: bool) -> str:
    parts = [f"M {fmt(points[0][0])} {fmt(points[0][1])}"]
    for x, y in points[1:]:
        parts.append(f"L {fmt(x)} {fmt(y)}")
    if closed:
        parts.append("Z")
    return " ".join(parts)


def svg_doc(viewbox: tuple[float, float, float, float], body: str, title: str) -> str:
    vb = " ".join(fmt(v) for v in viewbox)
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="{vb}" '
        f'role="img" aria-label="{title}">\n'
        f'<title>{title}</title>\n'
        f'{body}\n'
        f'</svg>\n'
    )


# ---------- Isometric brick ----------
def iso_block(W: int, D: int, H_units: int, palette: Palette,
              show_studs: bool = True, slope: bool = False,
              mirror: bool = False) -> str:
    """Emit a clean iso brick. H_units of 1 = plate, 3 = standard brick.

    Anchor: rear corner of top face placed at origin. When ``mirror`` is True,
    the block faces the opposite direction (iso-right axis flipped) and the
    wall shades are swapped so the sun stays in the upper left.
    """
    if W < 1 or D < 1 or H_units < 1:
        raise ValueError("W, D, H_units must all be >= 1")

    H = H_units * PLATE_H
    p = palette
    sx = -1 if mirror else 1

    # Mirroring flips the geometric right/left walls visually. Swap the shades
    # so the darker wall always appears on the visual right.
    right_fill = p.left if mirror else p.right
    left_fill = p.right if mirror else p.left

    # Top-face corners (rear -> right -> front -> left)
    rear = (0, 0)
    right = (sx * W * U, W * S)
    front = (sx * (W * U - D * U), W * S + D * S)
    left = (sx * (-D * U), D * S)

    # Bottom-face corners (drop by H)
    right_b = (right[0], right[1] + H)
    front_b = (front[0], front[1] + H)
    left_b = (left[0], left[1] + H)

    parts: list[str] = []

    if slope:
        # Slope descends along the D (long) axis with the stud at the rear end.
        # Real-world studded-brick proportions:
        #   1x2 slope is 45 degrees (drops all the way to the ground at the front)
        #   1x3 and longer are ~33 degrees (1 plate-height front wall remains)
        slope_drop = H if D <= 2 else H * 2 / 3
        left_top = (left[0], left[1] + slope_drop)
        front_top = (front[0], front[1] + slope_drop)

        # Left-front wall is a reduced-height slab; collapses to a line on a
        # 45-degree slope, so skip it then.
        if slope_drop < H:
            parts.append(poly([front_top, left_top, left_b, front_b], left_fill, p.outline))
        # Right-front wall: top edge follows the slope (full H at rear, tapering).
        parts.append(poly([front_top, right, right_b, front_b], right_fill, p.outline))
        # Sloped top face.
        parts.append(poly([rear, right, front_top, left_top], p.top, p.outline))

        if show_studs:
            for i in range(W):
                cx = rear[0] + sx * ((i + 0.5) * U + 0.5 * (-U))
                cy = rear[1] + (i + 0.5) * S + 0.5 * S
                _stud(parts, cx, cy, p)
        return "\n".join(parts)

    # Standard brick
    parts.append(poly([front, left, left_b, front_b], left_fill, p.outline))
    parts.append(poly([front, right, right_b, front_b], right_fill, p.outline))
    parts.append(poly([rear, right, front, left], p.top, p.outline))

    if show_studs:
        for i in range(W):
            for j in range(D):
                cx = rear[0] + sx * ((i + 0.5) * U + (j + 0.5) * (-U))
                cy = rear[1] + (i + 0.5) * S + (j + 0.5) * S
                _stud(parts, cx, cy, p)

    return "\n".join(parts)


def _stud(parts: list[str], cx: float, cy: float, p: Palette) -> None:
    """Append a raised cylindrical stud: bottom ellipse, side wall, top cap."""
    rx, ry, h = STUD_RX, STUD_RY, STUD_RAISE
    parts.append(ellipse(cx, cy, rx, ry, p.stud_side, p.outline))
    parts.append(
        f'<rect x="{fmt(cx - rx)}" y="{fmt(cy - h)}" '
        f'width="{fmt(2 * rx)}" height="{fmt(h)}" '
        f'fill="{p.stud_side}"/>'
    )
    parts.append(
        f'<path d="M {fmt(cx - rx)} {fmt(cy - h)} L {fmt(cx - rx)} {fmt(cy)} '
        f'M {fmt(cx + rx)} {fmt(cy - h)} L {fmt(cx + rx)} {fmt(cy)}" '
        f'fill="none" stroke="{p.outline}" stroke-width="{STROKE_W}"/>'
    )
    parts.append(ellipse(cx, cy - h, rx, ry, p.stud_top, p.outline))


def iso_viewbox(W: int, D: int, H_units: int,
                mirror: bool = False, pad: int = 4) -> tuple[float, float, float, float]:
    """Compute snug viewBox for iso block anchored at origin."""
    H = H_units * PLATE_H
    if mirror:
        min_x = -W * U
        max_x = D * U
    else:
        min_x = -D * U
        max_x = W * U
    min_y = -STUD_RAISE - STUD_RY
    max_y = W * S + D * S + H  # front corner drops by H
    return (min_x - pad, min_y - pad, (max_x - min_x) + 2 * pad, (max_y - min_y) + 2 * pad)


def _strip_bounds(W: int, D: int) -> tuple[float, float, float, float, float]:
    """Geometry shared by `strip_iso_block` and `strip_iso_viewbox`.

    Returns (slab_left, slab_right, slab_top_y, para_bot_y, front_bot_y).
    """
    slab_left = -2 * U
    slab_right = W * U
    slab_top_y = S - STUD_RY - STUD_RAISE - 2
    para_bot_y = D * S + STUD_RY + 2
    front_bot_y = para_bot_y + PLATE_H
    return slab_left, slab_right, slab_top_y, para_bot_y, front_bot_y


def strip_iso_block(W: int, D: int, palette: Palette) -> str:
    """Flat horizontal "horizon" slab with iso-lattice studs.

    Studs sit on the standard iso lattice (a*U, b*S where a+b is odd), the same
    grid existing iso bricks use, so a brick placed on the slab snaps cleanly:
    its front-most stud cell aligns with a slab stud and rear cells lift into
    the air at the iso angle. The slab silhouette is a thin horizontal
    rectangle — top face flat, front face flat — cropped to D visible stud rows
    rather than the full iso parallelogram height.
    """
    p = palette
    slab_left, slab_right, slab_top_y, para_bot_y, front_bot_y = _strip_bounds(W, D)

    parts = [
        poly(
            [(slab_left, para_bot_y), (slab_right, para_bot_y),
             (slab_right, front_bot_y), (slab_left, front_bot_y)],
            p.left, p.outline,
        ),
        poly(
            [(slab_left, slab_top_y), (slab_right, slab_top_y),
             (slab_right, para_bot_y), (slab_left, para_bot_y)],
            p.top, p.outline,
        ),
    ]
    for b in range(1, D + 1):
        for a in range(-2, W + 1):
            if (a + b) % 2 == 0:
                continue
            cx = a * U
            cy = b * S
            if cx < slab_left or cx > slab_right:
                continue
            _stud(parts, cx, cy, p)
    return "\n".join(parts)


def strip_iso_viewbox(W: int, D: int,
                      pad: int = 4) -> tuple[float, float, float, float]:
    slab_left, slab_right, slab_top_y, _, front_bot_y = _strip_bounds(W, D)
    return (
        slab_left - pad,
        slab_top_y - pad,
        (slab_right - slab_left) + 2 * pad,
        (front_bot_y - slab_top_y) + 2 * pad,
    )


# ---------- Broken (snapped-in-half) brick ----------
# Each half is clipped from the full iso body, rotated about the break
# midpoint so the snap end lifts (outer end drops), then translated outward
# along the D axis to open a small gap. Visual parameters (angle, gap,
# amplitude) are tuned once at module scope — they produce a consistent comic
# snap across all sizes.
_BROKEN_ANGLE_DEG = 7.0
_BROKEN_GAP = 5.0
_BROKEN_AMPLITUDE = 7.0
_BROKEN_STROKE_W = 1.4
_BROKEN_VB_PAD = 6


def broken_iso_body(
    W: int, D: int, H_units: int, palette: Palette,
    show_studs: bool = True, mirror: bool = False,
) -> tuple[str, tuple[float, float, float, float]]:
    """Snap a brick in half at D // 2 with a comic-style jagged edge.

    Returns (inner SVG content, tight viewBox). Caller wraps the content
    with `svg_doc`. For mirror=True, x-coords and rotation/translation signs
    flip so the snap reads the same relative to the viewer.
    """
    if D < 4 or D % 2 != 0:
        raise ValueError(f"broken brick needs even D >= 4, got {D}")
    d_break = D // 2
    sx = -1 if mirror else 1
    H = H_units * PLATE_H
    p = palette

    # Break endpoints on screen: top-face W=0 end, top-face W=W end, right-
    # wall bottom at W=W.
    top_start = (sx * -d_break * U, d_break * S)
    top_end = (sx * (W - d_break) * U, W * S + d_break * S)
    wall_bot = (top_end[0], top_end[1] + H)

    top_len = ((top_end[0] - top_start[0]) ** 2
               + (top_end[1] - top_start[1]) ** 2) ** 0.5
    # ~1 peak per 12 units keeps the zigzag density consistent across sizes.
    n_top = max(5, round(top_len / 12))
    n_wall = max(3, round(H / 10))
    # Cap wall amplitude proportionally to H so plate walls (H=12) don't get
    # zigzags that swing beyond the wall itself.
    wall_amp = min(_BROKEN_AMPLITUDE, H / 3)

    break_poly = (
        _zigzag(top_start, top_end, n_top, _BROKEN_AMPLITUDE)
        + _zigzag(top_end, wall_bot, n_wall, wall_amp)[1:]
    )

    # Silhouette corners of each half, used to close the clip polygons.
    rear = (0, 0)
    right_top = (sx * W * U, W * S)
    right_bot = (right_top[0], right_top[1] + H)
    front_top = (sx * (W - D) * U, W * S + D * S)
    front_bot = (front_top[0], front_top[1] + H)
    left_top = (sx * -D * U, D * S)
    left_bot = (left_top[0], left_top[1] + H)

    clip_a = break_poly + [right_bot, right_top, rear]
    clip_b = break_poly + [front_bot, left_bot, left_top]

    # Rotation pivot: midpoint of the straight break line from top_start to
    # wall_bot. Each half rotates so its outer end drops and the snap end
    # lifts. For mirror, the side swap requires inverted angle signs.
    mx = (top_start[0] + wall_bot[0]) / 2
    my = (top_start[1] + wall_bot[1]) / 2
    angle_a = sx * _BROKEN_ANGLE_DEG
    angle_b = -sx * _BROKEN_ANGLE_DEG

    d_len = (U * U + S * S) ** 0.5
    # D-axis unit vector on screen is (-sx*U, S)/d_len. Half-A shifts along
    # -D (toward rear), half-B along +D (toward front).
    tA = (sx * U * _BROKEN_GAP / d_len, -S * _BROKEN_GAP / d_len)
    tB = (-sx * U * _BROKEN_GAP / d_len, S * _BROKEN_GAP / d_len)

    body = iso_block(W, D, H_units, palette, show_studs=show_studs,
                     slope=False, mirror=mirror)

    defs = (
        '<defs>'
        f'<clipPath id="halfA"><path d="{_polyline_path(clip_a, True)}"/></clipPath>'
        f'<clipPath id="halfB"><path d="{_polyline_path(clip_b, True)}"/></clipPath>'
        '</defs>'
    )
    break_d = _polyline_path(break_poly, False)
    break_stroke = (
        f'<path d="{break_d}" fill="none" stroke="{p.outline}" '
        f'stroke-width="{_BROKEN_STROKE_W}" '
        'stroke-linejoin="round" stroke-linecap="round"/>'
    )

    def _tr(tx: float, ty: float, angle: float) -> str:
        return (f'translate({fmt(tx)} {fmt(ty)}) '
                f'rotate({fmt(angle)} {fmt(mx)} {fmt(my)})')

    content = (
        f'{defs}\n'
        f'<g transform="{_tr(tA[0], tA[1], angle_a)}">\n'
        f'<g clip-path="url(#halfA)">\n{body}\n</g>\n'
        f'{break_stroke}\n'
        f'</g>\n'
        f'<g transform="{_tr(tB[0], tB[1], angle_b)}">\n'
        f'<g clip-path="url(#halfB)">\n{body}\n</g>\n'
        f'{break_stroke}\n'
        f'</g>'
    )

    # Tight viewBox from the transformed half-silhouettes. Stud tops extend
    # above the top face by STUD_RAISE + STUD_RY, added to y_min.
    corners_a = [rear, right_top, right_bot, wall_bot, top_start]
    corners_b = [front_top, front_bot, left_bot, left_top, wall_bot, top_start]

    def _apply(pt: tuple[float, float], angle: float,
               tx: float, ty: float) -> tuple[float, float]:
        a = math.radians(angle)
        c, s = math.cos(a), math.sin(a)
        dx, dy = pt[0] - mx, pt[1] - my
        return (mx + dx * c - dy * s + tx, my + dx * s + dy * c + ty)

    pts = [_apply(pt, angle_a, tA[0], tA[1]) for pt in corners_a]
    pts += [_apply(pt, angle_b, tB[0], tB[1]) for pt in corners_b]
    xs = [q[0] for q in pts]
    ys = [q[1] for q in pts]
    x_min, x_max = min(xs), max(xs)
    y_min = min(ys) - (STUD_RAISE + STUD_RY)
    y_max = max(ys)
    pad = _BROKEN_VB_PAD
    vb = (x_min - pad, y_min - pad,
          (x_max - x_min) + 2 * pad, (y_max - y_min) + 2 * pad)
    return content, vb


def _centre_viewbox(tight_vb: tuple[float, float, float, float],
                    canvas_w: float, canvas_h: float) -> tuple[float, float, float, float]:
    """Return a canvas_w x canvas_h viewBox centred on the piece's bounding box.

    Used for uniform-canvas output so every PNG in a perspective has identical
    pixel dimensions — auto-fit in slide tools then applies the same scale to
    all of them and relative sizes stay correct.
    """
    tx, ty, tw, th = tight_vb
    cx = tx + tw / 2
    cy = ty + th / 2
    return (cx - canvas_w / 2, cy - canvas_h / 2, canvas_w, canvas_h)


# ---------- Top-down view ----------
def top_block(W: int, D: int, palette: Palette, show_studs: bool = True) -> str:
    p = palette
    cell = U
    body_w = W * cell
    body_h = D * cell
    parts: list[str] = []
    parts.append(rect(0, 0, body_w, body_h, p.top, p.outline, rx=2))
    if show_studs:
        for i in range(W):
            for j in range(D):
                cx = (i + 0.5) * cell
                cy = (j + 0.5) * cell
                parts.append(circle(cx, cy, cell * 0.3, p.stud_top, p.outline))
    return "\n".join(parts)


def top_viewbox(W: int, D: int, pad: int = 4) -> tuple[float, float, float, float]:
    return (-pad, -pad, W * U + 2 * pad, D * U + 2 * pad)


# ---------- Side view ----------
def side_block(W: int, H_units: int, palette: Palette, show_studs: bool = True) -> str:
    """Front-elevation: width = W studs, height = H plate-units, plus stud bumps."""
    p = palette
    body_w = W * U
    body_h = H_units * PLATE_H
    stud_bump = PLATE_H  # stud height on side view = 1 plate-unit
    stud_w = U * 0.6
    parts: list[str] = []
    parts.append(rect(0, stud_bump, body_w, body_h, p.top, p.outline, rx=1))
    parts.append(rect(0, stud_bump + body_h - 2, body_w, 2, p.right, "none"))
    if show_studs:
        for i in range(W):
            cx = (i + 0.5) * U
            parts.append(rect(cx - stud_w / 2, 0, stud_w, stud_bump + 1.5,
                              p.stud_top, p.outline, rx=1.5))
    return "\n".join(parts)


def side_viewbox(W: int, H_units: int, pad: int = 4) -> tuple[float, float, float, float]:
    return (-pad, -pad, W * U + 2 * pad, H_units * PLATE_H + PLATE_H + 2 * pad)


# ---------- Catalogue ----------
@dataclass(frozen=True)
class BlockSpec:
    slug: str       # filename slug
    W: int
    D: int
    H_units: int    # 1 = plate, 3 = brick, 6 = tall brick
    show_studs: bool = True
    slope: bool = False
    broken: bool = False  # render as snapped-in-half comic-style variant
    strip: bool = False   # flat horizontal "horizon" baseplate (iso view only)


STANDARD_BRICKS = [
    BlockSpec("brick-1x1", 1, 1, 3),
    BlockSpec("brick-1x2", 1, 2, 3),
    BlockSpec("brick-1x3", 1, 3, 3),
    BlockSpec("brick-1x4", 1, 4, 3),
    BlockSpec("brick-1x6", 1, 6, 3),
    BlockSpec("brick-1x8", 1, 8, 3),
    BlockSpec("brick-2x2", 2, 2, 3),
    BlockSpec("brick-2x3", 2, 3, 3),
    BlockSpec("brick-2x4", 2, 4, 3),
    BlockSpec("brick-2x6", 2, 6, 3),
    BlockSpec("brick-2x8", 2, 8, 3),
    BlockSpec("brick-2x10", 2, 10, 3),
    BlockSpec("brick-4x4", 4, 4, 3),
    BlockSpec("brick-4x6", 4, 6, 3),
]

PLATES = [
    BlockSpec("plate-1x1", 1, 1, 1),
    BlockSpec("plate-1x2", 1, 2, 1),
    BlockSpec("plate-1x4", 1, 4, 1),
    BlockSpec("plate-1x6", 1, 6, 1),
    BlockSpec("plate-1x8", 1, 8, 1),
    BlockSpec("plate-2x2", 2, 2, 1),
    BlockSpec("plate-2x4", 2, 4, 1),
    BlockSpec("plate-2x6", 2, 6, 1),
    BlockSpec("plate-2x8", 2, 8, 1),
    BlockSpec("plate-4x4", 4, 4, 1),
    BlockSpec("plate-4x6", 4, 6, 1),
    # Baseplates — the bigger single-layer sheets you build on top of.
    BlockSpec("plate-4x8", 4, 8, 1),
    BlockSpec("plate-4x10", 4, 10, 1),
    BlockSpec("plate-6x6", 6, 6, 1),
    BlockSpec("plate-8x8", 8, 8, 1),
    BlockSpec("plate-8x16", 8, 16, 1),
    BlockSpec("plate-16x16", 16, 16, 1),
    BlockSpec("plate-16x32", 16, 32, 1),
]

TILES = [
    BlockSpec("tile-1x1", 1, 1, 1, show_studs=False),
    BlockSpec("tile-1x2", 1, 2, 1, show_studs=False),
    BlockSpec("tile-1x4", 1, 4, 1, show_studs=False),
    BlockSpec("tile-2x2", 2, 2, 1, show_studs=False),
    BlockSpec("tile-2x4", 2, 4, 1, show_studs=False),
]

TALL = [
    BlockSpec("tall-1x2", 1, 2, 6),
    BlockSpec("tall-2x2", 2, 2, 6),
]

SLOPES = [
    BlockSpec("slope-1x2", 1, 2, 3, slope=True),
    BlockSpec("slope-1x3", 1, 3, 3, slope=True),
]

# Comic-book "snapped in half" variants. Break plane is at D // 2, so D must
# be even and >= 4 for a useful split. Iso view only — top/side views can't
# convey the snap. See `broken_iso_body` for geometry. Mix of bricks and
# plates: small snapped bricks read well, wider pieces only really work as
# plates because the wall zigzag on a 36px brick wall would dominate.
BROKEN_PIECES = [
    BlockSpec("broken-brick-2x4",  2, 4,  3, broken=True),
    BlockSpec("broken-brick-2x8",  2, 8,  3, broken=True),
    BlockSpec("broken-plate-2x10", 2, 10, 1, broken=True),
    BlockSpec("broken-plate-4x10", 4, 10, 1, broken=True),
]

# Flat horizontal "horizon" baseplates. Iso view only — designed to sit at the
# bottom of a slide with existing iso bricks placed on top. Studs follow the
# standard iso lattice (a*U, b*S where a+b odd) so any iso brick docks on the
# same grid. `D` here means visible stud-rows in screen depth, not iso cells —
# the slab silhouette is forced to a thin horizontal rectangle. See
# `strip_iso_block` for the geometry.
STRIPS = [
    BlockSpec("strip-32x2", 32, 2, 1, strip=True),
    BlockSpec("strip-32x4", 32, 4, 1, strip=True),
]


def render_iso_svg(
    spec: BlockSpec, palette: Palette, mirror: bool = False,
    canvas: tuple[float, float] | None = None,
) -> str:
    if spec.strip:
        body = strip_iso_block(spec.W, spec.D, palette)
        vb = strip_iso_viewbox(spec.W, spec.D)
        kind = "strip"
    elif spec.broken:
        body, vb = broken_iso_body(spec.W, spec.D, spec.H_units, palette,
                                   show_studs=spec.show_studs, mirror=mirror)
        kind = "broken"
    else:
        body = iso_block(spec.W, spec.D, spec.H_units, palette,
                         show_studs=spec.show_studs, slope=spec.slope,
                         mirror=mirror)
        vb = iso_viewbox(spec.W, spec.D, spec.H_units, mirror=mirror)
        kind = "isometric"
    if canvas is not None:
        vb = _centre_viewbox(vb, *canvas)
    facing = " mirror" if mirror else ""
    title = f"{spec.slug} {palette.name} ({kind}{facing})"
    return svg_doc(vb, body, title)


def render_top_svg(
    spec: BlockSpec, palette: Palette,
    canvas: tuple[float, float] | None = None,
) -> str:
    body = top_block(spec.W, spec.D, palette, show_studs=spec.show_studs)
    vb = top_viewbox(spec.W, spec.D)
    if canvas is not None:
        vb = _centre_viewbox(vb, *canvas)
    title = f"{spec.slug} {palette.name} (top-down)"
    return svg_doc(vb, body, title)


def render_side_svg(
    spec: BlockSpec, palette: Palette,
    canvas: tuple[float, float] | None = None,
) -> str:
    body = side_block(spec.W, spec.H_units, palette, show_studs=spec.show_studs)
    vb = side_viewbox(spec.W, spec.H_units)
    if canvas is not None:
        vb = _centre_viewbox(vb, *canvas)
    title = f"{spec.slug} {palette.name} (side)"
    return svg_doc(vb, body, title)


def max_iso_canvas(specs: list[BlockSpec], pad: int = 4) -> tuple[float, float]:
    """Canvas size that fits every iso piece in the list (both facings).

    Iso mirror flips the x-range but preserves width, so one canvas works for
    both. Scale is driven by the largest `(W + D)` piece plus the largest `H`
    among others (e.g. TALL extends y only).
    """
    w = max((s.W + s.D) * U for s in specs) + 2 * pad
    h_content = max((s.W + s.D) * S + s.H_units * PLATE_H for s in specs)
    h = h_content + STUD_RAISE + STUD_RY + 2 * pad
    return w, h


def max_top_canvas(specs: list[BlockSpec], pad: int = 4) -> tuple[float, float]:
    w = max(s.W * U for s in specs) + 2 * pad
    h = max(s.D * U for s in specs) + 2 * pad
    return w, h


def max_side_canvas(specs: list[BlockSpec], pad: int = 4) -> tuple[float, float]:
    w = max(s.W * U for s in specs) + 2 * pad
    h = max(s.H_units * PLATE_H for s in specs) + PLATE_H + 2 * pad
    return w, h


def dedupe_for_side(specs: list[BlockSpec]) -> list[BlockSpec]:
    """Side elevation depends only on (W, H_units, show_studs).

    Two bricks with the same W and H render identical side views regardless
    of D. Keep one canonical representative (smallest D) per tuple.
    """
    best: dict[tuple[int, int, bool], BlockSpec] = {}
    for s in specs:
        key = (s.W, s.H_units, s.show_studs)
        if key not in best or s.D < best[key].D:
            best[key] = s
    return list(best.values())


# ---------- Build ----------
def _palette_allows(spec: BlockSpec, palette: Palette) -> bool:
    """Accent palettes (very light tones) only render on curated specs."""
    if palette.accent:
        return spec.slug in ACCENT_SLUGS
    return True


def _resolve_palette_sets(spec: str) -> list[str]:
    """Turn a comma-separated --palettes value into a list of palette-set names."""
    names = [p.strip() for p in spec.split(",") if p.strip()]
    if "all" in names:
        return list(PALETTE_SETS)
    for n in names:
        if n not in PALETTE_SETS:
            raise SystemExit(
                f"Unknown palette set: {n!r}. "
                f"Available: {', '.join(PALETTE_SETS)}, all"
            )
    return names


def _try_import_cairosvg():
    """Return the cairosvg module if available, else None.

    cairosvg is optional — pure SVG output has no external deps. Only PNG
    rasterisation needs it.
    """
    try:
        import cairosvg  # type: ignore[import-not-found]
        return cairosvg
    except ImportError:
        return None


def _rasterise_one(job: tuple[Path, Path, int]) -> None:
    """Worker: render one SVG to PNG. Each process imports cairosvg itself."""
    svg_path, png_path, output_width = job
    import cairosvg  # type: ignore[import-not-found]
    png_path.parent.mkdir(parents=True, exist_ok=True)
    cairosvg.svg2png(
        bytestring=svg_path.read_bytes(),
        write_to=str(png_path),
        output_width=output_width,
    )


def _rasterise_dir(svg_dir: Path, png_dir: Path, scale: int, min_width: int, workers: int) -> int:
    """Mirror svg_dir as PNGs into png_dir. Returns count rendered.

    Jobs are dispatched to a ProcessPoolExecutor when there's enough work to
    amortise pool startup; small runs stay serial.
    """
    if png_dir.exists():
        shutil.rmtree(png_dir)
    png_dir.mkdir(parents=True, exist_ok=True)

    jobs: list[tuple[Path, Path, int]] = []
    for svg in svg_dir.rglob("*.svg"):
        text = svg.read_text()
        m = VB_RE.search(text)
        if not m:
            continue
        vb_w = float(m.group(1).split()[2])
        out_w = max(min_width, int(round(vb_w * scale)))
        rel = svg.relative_to(svg_dir).with_suffix(".png")
        jobs.append((svg, png_dir / rel, out_w))

    if workers <= 1 or len(jobs) < 8:
        for job in jobs:
            _rasterise_one(job)
        return len(jobs)

    with ProcessPoolExecutor(max_workers=workers) as pool:
        # Consume the map iterator so worker exceptions surface here.
        for _ in pool.map(_rasterise_one, jobs):
            pass
    return len(jobs)


# ---------- PPTX generation ----------
# Points (pt) are PowerPoint's natural unit; 72 pt = 1 inch, 12700 EMU = 1 pt.
# 16:9 widescreen default: 13.333" x 7.5" = 960 x 540 pt.
SLIDE_W_PT = 960.0
SLIDE_H_PT = 540.0
PPTX_MARGIN_PT = 8.0
PPTX_GAP_PT = 6.0
PPTX_DENSITY_DEFAULT = 0.25  # pt per SVG user unit; tune with --pptx-density


@dataclass(frozen=True)
class _Placement:
    path: Path
    x: float  # pt from slide top-left
    y: float
    w: float
    h: float


# Layout order for the PPTX. Baseplates are the *big* plates (W*D >= 32);
# they sort to the tail of the deck so the small bricks/slopes/tall/tiles/
# small-plates don't end up orphaned next to giants they can't line up with.
# Iso non-mirror separates from iso-mirror so same-facing pieces cluster.
_PPTX_CATEGORY_ORDER = {"brick": 0, "broken": 1, "slope": 2, "tall": 3, "tile": 4, "plate": 5, "strip": 6}
_PPTX_VIEW_ORDER = {"iso": 0, "side": 1, "top": 2}
_PPTX_SIZE_RE = re.compile(r"-(\d+)x(\d+)-")
_BASEPLATE_MIN_AREA = 32  # W*D >= 32 matches the "Baseplates" comment in PLATES


def _pptx_parse(name: str) -> tuple[str, int, int]:
    """Extract (category, W, D) from a PNG stem like 'plate-4x8-ocean-iso'.
    Returns ('', 0, 0) if the slug doesn't match the expected pattern.
    """
    category = name.split("-", 1)[0]
    m = _PPTX_SIZE_RE.search(name)
    if not m:
        return category, 0, 0
    return category, int(m.group(1)), int(m.group(2))


def _is_baseplate(name: str) -> bool:
    category, w, d = _pptx_parse(name)
    if category == "strip":
        return True
    return category == "plate" and w * d >= _BASEPLATE_MIN_AREA


def _pptx_sort_key(png_path: Path, png_set_root: Path) -> tuple:
    """Sort key: (baseplate?, view, mirror, category, w*d, w, d, name).

    Primary key is is_baseplate so all small pieces come first across every
    view, then all baseplates. Secondary sort matches the within-view order
    from before: iso non-mirror, iso mirror, side, top; then by category.
    """
    rel = png_path.relative_to(png_set_root)
    view_dir = rel.parts[0]
    name = rel.stem
    is_mirror = name.endswith("-iso-mirror")
    category, w, d = _pptx_parse(name)
    return (
        1 if _is_baseplate(name) else 0,
        _PPTX_VIEW_ORDER.get(view_dir, 99),
        1 if is_mirror else 0,
        _PPTX_CATEGORY_ORDER.get(category, 99),
        w * d,
        w,
        d,
        name,
    )


def _try_import_pptx():
    """Return the pptx package if available, else None."""
    try:
        import pptx  # type: ignore[import-not-found]
        return pptx
    except ImportError:
        return None


def _parse_viewbox_wh(svg_path: Path) -> tuple[float, float]:
    """Return (width, height) in SVG user units from the viewBox attribute."""
    m = VB_RE.search(svg_path.read_text())
    if not m:
        raise ValueError(f"No viewBox in {svg_path}")
    parts = m.group(1).split()
    return float(parts[2]), float(parts[3])


def _pack_slides(
    items: list[tuple[Path, float, float]],
    points_per_unit: float,
) -> list[list[_Placement]]:
    """Lay items out one piece-family per row for catalogue-style scanning.

    A "family" is (view, is_mirror, category, W, D). All palette variants of,
    say, brick-2x4 share a family and land on a contiguous block so you can
    scan across to compare colours. A family never splits across slides — if
    the whole family won't fit on the current slide, we start a new one. Only
    exception: families that are physically too large to fit on any single
    slide (very dense decks with big baseplates) are split into chunks, each
    chunk on its own slide.

    Slide breaks also fire on (is_baseplate, is_iso_mirror) transitions so
    iso-right, iso-left, and giant baseplates each get their own slides.
    """
    usable_w = SLIDE_W_PT - 2 * PPTX_MARGIN_PT
    usable_h = SLIDE_H_PT - 2 * PPTX_MARGIN_PT
    gap = PPTX_GAP_PT
    margin = PPTX_MARGIN_PT

    # First pass: clamp oversized pieces to fit the surface.
    prepared: list[tuple[Path, float, float]] = []
    for path, vb_w, vb_h in items:
        w = vb_w * points_per_unit
        h = vb_h * points_per_unit
        if w > usable_w:
            s = usable_w / w
            w *= s
            h *= s
        if h > usable_h:
            s = usable_h / h
            w *= s
            h *= s
        prepared.append((path, w, h))

    # Second pass: group consecutive items by family. The input is pre-sorted
    # so all same-family items are already contiguous.
    families: list[tuple[tuple, list[tuple[Path, float, float]]]] = []
    cur_key: tuple | None = None
    cur_items: list[tuple[Path, float, float]] = []
    for path, w, h in prepared:
        name = path.stem
        is_baseplate = _is_baseplate(name)
        is_iso_mirror = name.endswith("-iso-mirror")
        category, piece_w, piece_d = _pptx_parse(name)
        view = path.parent.name  # 'iso', 'side', 'top'
        key = (is_baseplate, view, is_iso_mirror, category, piece_w, piece_d)
        if cur_key is not None and key != cur_key:
            families.append((cur_key, cur_items))
            cur_items = []
        cur_key = key
        cur_items.append((path, w, h))
    if cur_items and cur_key is not None:
        families.append((cur_key, cur_items))

    # Third pass: pack families onto slides. Single-row families share rows
    # with their neighbours when the combined widths fit; multi-row families
    # claim their own vertical band.
    slides: list[list[_Placement]] = []
    slide_items: list[_Placement] = []
    x = 0.0  # row cursor (surface-local; margin added at placement)
    y = 0.0  # row top (surface-local)
    row_h = 0.0
    prev_break_key: tuple[bool, bool] | None = None

    def flush() -> None:
        nonlocal slide_items, x, y, row_h
        if slide_items:
            slides.append(slide_items)
        slide_items = []
        x = 0.0
        y = 0.0
        row_h = 0.0

    for family_key, family_items in families:
        is_baseplate = family_key[0]
        is_iso_mirror = family_key[2]
        break_key = (is_baseplate, is_iso_mirror)

        # Fresh slide on quarantine boundaries (iso-right ↔ iso-left ↔
        # non-iso, and baseplates ↔ non-baseplates).
        if prev_break_key is not None and break_key != prev_break_key:
            flush()
        prev_break_key = break_key

        # Every item in a family has identical rendered dimensions (same W,
        # D, view, mirror, category), so the first item's size is the family
        # size.
        first_w = family_items[0][1]
        first_h = family_items[0][2]
        items_per_row = max(1, int((usable_w + gap) // (first_w + gap)))
        max_rows = max(1, int((usable_h + gap) // (first_h + gap)))
        max_items_per_slide = max_rows * items_per_row

        # Split oversized families (rare) into slide-sized chunks.
        i = 0
        while i < len(family_items):
            chunk = family_items[i:i + max_items_per_slide]
            i += len(chunk)
            chunk_n = len(chunk)

            # Single-row chunks try to share the current row with previous
            # families; multi-row chunks claim their own rows.
            if chunk_n <= items_per_row:
                chunk_width = chunk_n * first_w + (chunk_n - 1) * gap
                if (slide_items
                        and x + chunk_width <= usable_w
                        and y + first_h <= usable_h):
                    for path, w, h in chunk:
                        slide_items.append(_Placement(
                            path, x + margin, y + margin, w, h,
                        ))
                        x += w + gap
                    if first_h > row_h:
                        row_h = first_h
                    continue

            # Wrap to a new row for the chunk.
            if x > 0:
                y += row_h + gap
                x = 0.0
                row_h = 0.0

            # Whole-chunk vertical check — if it won't fit below the current
            # cursor, bump to a new slide so the family stays contiguous.
            chunk_rows = (chunk_n + items_per_row - 1) // items_per_row
            chunk_height = chunk_rows * first_h + (chunk_rows - 1) * gap
            if slide_items and y + chunk_height > usable_h:
                flush()

            for path, w, h in chunk:
                # Wrap within the chunk when a row fills.
                if x > 0 and x + w > usable_w:
                    y += row_h + gap
                    x = 0.0
                    row_h = 0.0
                slide_items.append(_Placement(
                    path, x + margin, y + margin, w, h,
                ))
                x += w + gap
                if h > row_h:
                    row_h = h

    if slide_items:
        slides.append(slide_items)
    return slides


def _collect_pngs_for_set(
    set_name: str, svg_root: Path, png_root: Path,
) -> list[tuple[Path, float, float]]:
    """Return (png_path, vb_w, vb_h) for every PNG under png_root/set_name.

    Sort order (see `_pptx_sort_key`): iso non-mirror → iso mirror → side →
    top; within each view, by category (brick, slope, tall, tile, plate) then
    ascending size. Dimensions come from the SVG's viewBox (PNG pixel dims are
    a scaled derivative and not needed — slide dimensions are set explicitly).
    """
    png_set = png_root / set_name
    svg_set = svg_root / set_name
    pngs = list(png_set.rglob("*.png"))
    pngs.sort(key=lambda p: _pptx_sort_key(p, png_set))
    items: list[tuple[Path, float, float]] = []
    for png in pngs:
        rel = png.relative_to(png_set).with_suffix(".svg")
        svg = svg_set / rel
        if not svg.is_file():
            continue
        vb_w, vb_h = _parse_viewbox_wh(svg)
        items.append((png, vb_w, vb_h))
    return items


def _generate_pptx(
    svg_root: Path, png_root: Path, out_path: Path,
    set_names: list[str], points_per_unit: float,
) -> tuple[int, int, int]:
    """Write one PPTX per palette set. Returns (file_count, slide_count, image_count).

    The palette name is always inserted before the suffix, e.g. `Building Blocks.pptx`
    with palette set ['mantel'] writes `Building Blocks - mantel.pptx`.
    """
    if _try_import_pptx() is None:
        print("note: python-pptx not installed — skipping PPTX. "
              "Install with `pip install python-pptx` to enable.")
        return 0, 0, 0

    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.util import Emu, Pt  # type: ignore[import-not-found]

    file_count = 0
    total_slides = 0
    total_images = 0

    for set_name in set_names:
        items = _collect_pngs_for_set(set_name, svg_root, png_root)
        if not items:
            print(f"note: no PNGs found for palette set {set_name!r} — skipping.")
            continue

        target = out_path.with_name(
            f"{out_path.stem} - {set_name}{out_path.suffix}"
        )

        prs = Presentation()
        prs.slide_width = Emu(int(SLIDE_W_PT * 12700))
        prs.slide_height = Emu(int(SLIDE_H_PT * 12700))
        # Prefer the named Blank layout; fall back to the last layout which
        # is usually also a plain one.
        blank_layout = next(
            (layout for layout in prs.slide_layouts if layout.name == "Blank"),
            prs.slide_layouts[-1],
        )

        set_slides = 0
        set_images = 0
        for slide_placements in _pack_slides(items, points_per_unit):
            slide = prs.slides.add_slide(blank_layout)
            for pl in slide_placements:
                slide.shapes.add_picture(
                    str(pl.path), Pt(pl.x), Pt(pl.y), Pt(pl.w), Pt(pl.h),
                )
                set_images += 1
            set_slides += 1

        prs.save(str(target))
        file_count += 1
        total_slides += set_slides
        total_images += set_images
        print(f"  {target.name}: {set_images} images across {set_slides} slides")

    return file_count, total_slides, total_images


def _render_palette_set(
    set_name: str,
    set_palettes: dict[str, Palette],
    out_root: Path,
    all_specs: list[BlockSpec],
    flat_specs: list[BlockSpec],
    side_specs: list[BlockSpec],
    skip_mirror: bool,
    iso_canvas: tuple[float, float] | None,
    top_canvas: tuple[float, float] | None,
    side_canvas: tuple[float, float] | None,
) -> tuple[int, int, int]:
    """Write one palette set into out_root/<set_name>/{iso,top,side}/.

    Any existing content under the set's subdir is removed first so stale
    files from previous catalogue shapes don't linger.
    """
    set_out = out_root / set_name
    if set_out.exists():
        shutil.rmtree(set_out)
    iso_dir = set_out / "iso"
    top_dir = set_out / "top"
    side_dir = set_out / "side"
    for d in (iso_dir, top_dir, side_dir):
        d.mkdir(parents=True, exist_ok=True)

    iso_count = 0
    for spec in all_specs:
        for pname, palette in set_palettes.items():
            if not _palette_allows(spec, palette):
                continue
            (iso_dir / f"{spec.slug}-{pname}-iso.svg").write_text(
                render_iso_svg(spec, palette, mirror=False, canvas=iso_canvas)
            )
            iso_count += 1
            # Square blocks are symmetric under mirror, and strip slabs are
            # left-right symmetric by construction; skip the duplicate.
            if not skip_mirror and spec.W != spec.D and not spec.strip:
                (iso_dir / f"{spec.slug}-{pname}-iso-mirror.svg").write_text(
                    render_iso_svg(spec, palette, mirror=True, canvas=iso_canvas)
                )
                iso_count += 1

    top_count = 0
    for spec in flat_specs:
        for pname, palette in set_palettes.items():
            if palette.accent:
                continue
            (top_dir / f"{spec.slug}-{pname}-top.svg").write_text(
                render_top_svg(spec, palette, canvas=top_canvas)
            )
            top_count += 1

    side_count = 0
    for spec in side_specs:
        for pname, palette in set_palettes.items():
            if palette.accent:
                continue
            (side_dir / f"{spec.slug}-{pname}-side.svg").write_text(
                render_side_svg(spec, palette, canvas=side_canvas)
            )
            side_count += 1

    return iso_count, top_count, side_count


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Parametric isometric studded-brick generator."
    )
    parser.add_argument(
        "--out", type=Path,
        default=Path(__file__).resolve().parent / "svg",
        help="Output directory (default: ./svg next to this script)",
    )
    parser.add_argument(
        "--palettes", default=next(iter(PALETTE_SETS)),
        help=("Comma-separated palette sets. Choices: "
              + ", ".join(PALETTE_SETS) + f", all. Default: "
              f"{next(iter(PALETTE_SETS))!r} only — pass 'all' or a "
              "comma-separated list to render more."),
    )
    parser.add_argument(
        "--no-mirror", action="store_true",
        help="Skip the mirrored (left-facing) iso variant",
    )
    parser.add_argument(
        "--png-out", type=Path,
        default=Path(__file__).resolve().parent / "png",
        help="PNG output directory (default: ./png next to this script)",
    )
    parser.add_argument(
        "--no-png", action="store_true",
        help="Skip PNG rasterisation (SVG-only output)",
    )
    parser.add_argument(
        "--scale", type=int, default=4,
        help="PNG pixels per SVG unit (default: 4)",
    )
    parser.add_argument(
        "--min-width", type=int, default=200,
        help="Minimum PNG width in pixels (default: 200)",
    )
    parser.add_argument(
        "--workers", type=int, default=max(1, (os.cpu_count() or 2) - 1),
        help="Parallel PNG workers (default: cpu_count - 1)",
    )
    parser.add_argument(
        "--uniform-canvas", action="store_true",
        help="Pad every image in a perspective to the same pixel dimensions "
             "(piece centred, transparent padding). Slide tools that auto-fit "
             "imports then preserve relative scale across pieces. Default is "
             "tight canvas — smaller files, but relative sizes can get "
             "stretched by tool-side auto-fit.",
    )
    parser.add_argument(
        "--pptx", dest="pptx", action="store_true", default=True,
        help="Generate 'Building Blocks.pptx' with PNGs placed at correct relative "
             "scale (default on when PNGs are rendered). Sidesteps slide-tool "
             "auto-fit entirely because dimensions are set in absolute pt.",
    )
    parser.add_argument(
        "--no-pptx", dest="pptx", action="store_false",
        help="Skip PPTX generation.",
    )
    parser.add_argument(
        "--pptx-out", type=Path,
        default=Path(__file__).resolve().parent / "Building Blocks.pptx",
        help=("PPTX output path stem (default: ./Building Blocks.pptx). The "
              "palette name is always inserted before the suffix, e.g. "
              "'Building Blocks - mantel.pptx'."),
    )
    parser.add_argument(
        "--pptx-density", type=float, default=PPTX_DENSITY_DEFAULT,
        help=("Points per SVG user-unit — higher = bigger pieces, fewer per "
              f"slide (default: {PPTX_DENSITY_DEFAULT})"),
    )
    parser.add_argument(
        "--pptx-palettes", default=None,
        help=("Palette sets to include in the PPTX (comma-separated, or 'all'). "
              "Default: same as --palettes. Override to build a combined deck "
              "from palette sets already on disk without regenerating them."),
    )
    args = parser.parse_args()

    set_names = _resolve_palette_sets(args.palettes)
    out: Path = args.out
    out.mkdir(parents=True, exist_ok=True)

    all_specs = STANDARD_BRICKS + BROKEN_PIECES + PLATES + TILES + TALL + SLOPES + STRIPS
    flat_specs = [s for s in all_specs
                  if not s.slope and not s.broken and not s.strip and s.H_units <= 3]
    side_specs = dedupe_for_side(flat_specs)

    if args.uniform_canvas:
        iso_canvas = max_iso_canvas(all_specs)
        top_canvas = max_top_canvas(flat_specs)
        side_canvas = max_side_canvas(side_specs)
    else:
        iso_canvas = top_canvas = side_canvas = None

    # PNG step is optional; warn once up-front if requested but unavailable.
    cairosvg_mod = None
    if not args.no_png:
        cairosvg_mod = _try_import_cairosvg()
        if cairosvg_mod is None:
            print("note: cairosvg not installed — skipping PNG rendering. "
                  "Install with `pip install cairosvg` to enable.")

    svg_total = 0
    png_total = 0
    for set_name in set_names:
        iso_c, top_c, side_c = _render_palette_set(
            set_name, PALETTE_SETS[set_name],
            out, all_specs, flat_specs, side_specs,
            skip_mirror=args.no_mirror,
            iso_canvas=iso_canvas,
            top_canvas=top_canvas,
            side_canvas=side_canvas,
        )
        set_svg = iso_c + top_c + side_c
        svg_total += set_svg
        msg = f"  {set_name}: {set_svg} SVGs (iso={iso_c}, top={top_c}, side={side_c})"
        if cairosvg_mod is not None:
            set_png = _rasterise_dir(
                out / set_name, args.png_out / set_name,
                args.scale, args.min_width, args.workers,
            )
            png_total += set_png
            msg += f", {set_png} PNGs"
        print(msg)

    if cairosvg_mod is not None:
        print(f"Generated {svg_total} SVGs in {out} and {png_total} PNGs in {args.png_out}")
    else:
        print(f"Generated {svg_total} SVGs in {out}")

    if args.pptx:
        if args.no_png or cairosvg_mod is None:
            # The PPTX step reads PNGs from disk; no PNGs means nothing to place.
            print("note: skipping PPTX — needs PNGs "
                  "(re-run without --no-png, with cairosvg installed).")
        else:
            pptx_sets = _resolve_palette_sets(
                args.pptx_palettes if args.pptx_palettes else args.palettes,
            )
            file_count, slide_count, image_count = _generate_pptx(
                out, args.png_out, args.pptx_out,
                pptx_sets, args.pptx_density,
            )
            if file_count > 0:
                suffix = "" if file_count == 1 else f" across {file_count} PPTX files"
                print(f"Wrote {image_count} images across {slide_count} slides"
                      f"{suffix}")


if __name__ == "__main__":
    main()
